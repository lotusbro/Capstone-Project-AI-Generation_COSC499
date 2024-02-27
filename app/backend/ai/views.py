from django.shortcuts import render
from django.http import HttpResponse, JsonResponse, FileResponse, Http404
from django.views.decorators.csrf import csrf_exempt
from dotenv import load_dotenv
from PyPDF2 import PdfReader
from langchain.text_splitter import CharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_openai import OpenAIEmbeddings, ChatOpenAI
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser
from langchain_core.runnables import RunnablePassthrough
from datetime import datetime
import xml.etree.ElementTree as ET
from pptx.dml.color import RGBColor
from pptx import Presentation
import os

# Create your views here.

system_prompt = """You are an agent whose task is to generate text in order to generate a PowerPoint presentation. 

In order to do this, you need generate text in an XML style format.

You will be using preset slide layouts. Each of these slide layouts has a list of placeholders that you will populate with text.

Here is an example of the format, showing each of the 6 preset slide layouts:

<slides>
    <slide layout="title">
        <title>Slide Title</title>
        <subtitle>Slide Subtitle</subtitle>
    </slide>
    <slide layout="content">
        <title>Slide Title</title>
        <content>
			<b>Point 1</b>
			<b>Point 2</b>
		</content>
    </slide>
    <slide layout="header">
	    <title>Slide Title</title>
	    <text>Text</text>
    </slide>
    <slide layout="two">
	    <title>Slide Title</title>
	    <content>
			<b>Point 1</b>
			<b>Point 2</b>
		</content>
	    <content>
			<b>Point 1</b>
			<b>Point 2</b>
	    </content>
    </slide>
    <slide layout="comp">
	    <title>Slide Title</title>
	    <text>Left Side Title</text>
	    <content>
			<b>Point 1</b>
			<b>Point 2</b>
	    </content>
	    <text>Right Side Title</text>
	    <content>
			<b>Point 1</b>
			<b>Point 2</b>
	    </content>
    </slide>
    <slide layout="caption">
	    <title>Slide Title</title>
	    <content>
			<b>Point 1</b>
			<b>Point 2</b>
	    </content>
	    <text>Caption</text>
    </slide>
</slides>

Please note that you do not need to use every single layout preset, only use the preset you determine to be the most effective at displaying the information for a particular slide.

Here is a description of each of the slide layouts for you to determine their usefulness:
"title": A slide with a centered title and optional subtitle
"content": A slide with a title on top and a content body below
"header": Shorthand for section header. A slide used to signal a new section in the presentation, essentially the same as a title slide but styled differently
"two": Short for two-content, a slide with a title on top, and two side by side content bodies
"comp": Short for comparison, the same thing as a two content slide, but with a title for each side
"caption": A slide with a title and caption on the left side, and a main content body on the right

Please also note that a content placeholder displays text as bullet points for every new line, whereas title/subtitle/text do not, so you will need to wrap each bullet point with a `<b>` tag.

Here is an example of what you need to produce:
<slides>
    <slide layout="title">
        <title>Understanding Subject</title>
        <subtitle>Looking at Subject</subtitle>
    </slide>
    <slide layout="content">
        <title>What is Topic?</title>
        <content>
			<b>Point 1</b>
			<b>Point 2</b>
		</content>
    </slide>
    <slide layout="content">
        <title>Why is topic Important?</title>
        <content>
			<b>Helps us learn new things</b>
		</content>
    </slide>
    <slide layout="comp">
	    <title>Topic</title>
	    <text>What is Topic?</text>
	    <content>
			<b>Point 1</b>
			<b>Point 2</b>
	    </content>
	    <text>Why is it ___?</text>
	    <content>
			<b>Point 1</b>
			<b>Point 2</b>
	    </content>
    </slide>
    <slide layout="content">
        <title>What is Topic?</title>
        <content>
			<b>Point 1</b>
            <b>Example: </b>
		</content>
    </slide>
    <slide layout="title">
        <title>Recap: </title>
        <subtitle>Understanding</subtitle>
    </slide>
</slides>

Please leave out any string used in markdown e.g. ```xml```
"""

template = """{system}
Based on the following text parsed from {ctx}. Help adapt this into new materials for a presentation can be used to assist learning for other age groups.

{document}

Please use the original materials and convert them into slides that can be can be used to teach a grade {targetGrade} class. {prompt}."""

current_file_path = os.path.dirname(os.path.realpath(__file__))
GENERATEDCONTENT_DIRECTORY = os.path.join(current_file_path, 'generatedcontent')

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
LIGHTBLUE = RGBColor(135, 206, 235)
CREAM = RGBColor(255, 253, 208)
GREY = RGBColor(128, 128, 128)

def generate_filename(extension='.pptx', directory=GENERATEDCONTENT_DIRECTORY):
    timestamp = datetime.now().strftime('%Y%m%d%H%M%S')
    file_name = f"presentation_{timestamp}{extension}"
    
    # handle duplicate filename case
    counter = 1
    unique_file_name = file_name
    while os.path.exists(os.path.join(directory, unique_file_name)):
        unique_file_name = f"presentation_{timestamp}_{counter}{extension}"
        counter += 1
    return unique_file_name

def set_background_color(prs, color):
    for master in prs.slide_masters:
        background = master.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color  # Set color to ---- (needs a form value to change here)

def set_font(prs, type):
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = type  # Set font to ---- (needs a form value to change here)

def set_font_color(prs, color):
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = color  # Set font color ---- (needs a form value to change)

def generate_slides_from_XML(xml_string, bgcolor, fonttype, fontcolor):
    root = ET.fromstring(xml_string)
    prs = Presentation()
    set_background_color(prs, bgcolor)

    for slide in root.findall('slide'):
        if slide.get('layout') == 'title':
            title_slide_layout = prs.slide_layouts[0]
            title_slide = prs.slides.add_slide(title_slide_layout)
            shapes = title_slide.shapes
            for i, element in enumerate(slide):
                if element.tag == "content":
                    content_string = "" 
                    for bullet in element:
                        content_string += (bullet.text + '\n')
                    shapes.placeholders[i].text_frame.text = content_string
                else:
                    shapes.placeholders[i].text_frame.text = element.text
        elif slide.get('layout') == 'content':
            content_slide_layout = prs.slide_layouts[1]
            content_slide = prs.slides.add_slide(content_slide_layout)
            shapes = content_slide.shapes
            for i, element in enumerate(slide):
                if element.tag == "content":
                    content_string = "" 
                    for bullet in element:
                        content_string += (bullet.text + '\n')
                    shapes.placeholders[i].text_frame.text = content_string
                else:
                    shapes.placeholders[i].text_frame.text = element.text
        elif slide.get('layout') == 'header':
            header_slide_layout = prs.slide_layouts[2]
            header_slide = prs.slides.add_slide(header_slide_layout)
            shapes = header_slide.shapes
            for i, element in enumerate(slide):
                if element.tag == "content":
                    content_string = "" 
                    for bullet in element:
                        content_string += (bullet.text + '\n')
                    shapes.placeholders[i].text_frame.text = content_string
                else:
                    shapes.placeholders[i].text_frame.text = element.text
        elif slide.get('layout') == 'two':
            twocontent_slide_layout = prs.slide_layouts[3]
            twocontent_slide = prs.slides.add_slide(twocontent_slide_layout)
            shapes = twocontent_slide.shapes
            for i, element in enumerate(slide):
                if element.tag == "content":
                    content_string = "" 
                    for bullet in element:
                        content_string += (bullet.text + '\n')
                    shapes.placeholders[i].text_frame.text = content_string
                else:
                    shapes.placeholders[i].text_frame.text = element.text
        elif slide.get('layout') == 'comp':
            comparison_slide_layout = prs.slide_layouts[4]
            comparison_slide = prs.slides.add_slide(comparison_slide_layout)
            shapes = comparison_slide.shapes
            for i, element in enumerate(slide):
                if element.tag == "content":
                    content_string = "" 
                    for bullet in element:
                        content_string += (bullet.text + '\n')
                    shapes.placeholders[i].text_frame.text = content_string
                else:
                    shapes.placeholders[i].text_frame.text = element.text
        elif slide.get('layout') == 'caption':
            caption_slide_layout = prs.slide_layouts[7]
            caption_slide = prs.slides.add_slide(caption_slide_layout)
            shapes = caption_slide.shapes
            for i, element in enumerate(slide):
                if element.tag == "content":
                    content_string = "" 
                    for bullet in element:
                        content_string += (bullet.text + '\n')
                    shapes.placeholders[i].text_frame.text = content_string
                else:
                    shapes.placeholders[i].text_frame.text = element.text
    set_font(prs, fonttype)
    set_font_color(prs, fontcolor)
    file_name = generate_filename()
    file_path = os.path.join(GENERATEDCONTENT_DIRECTORY, file_name)
    prs.save(file_path)
    return file_name

@csrf_exempt
def ai(request):
    if (request.method == 'POST'):
        load_dotenv()
        uploaded_file = request.FILES.get('file')
        user_prompt = request.POST.get('prompt')
        ctx = request.POST.get('ctx')
        targetGrade = request.POST.get('targetGrade')
        bgcolor = request.POST.get('backgroundColor')
        fonttype = request.POST.get('fontType')
        fontcolor = request.POST.get('fontColor')

        # extract text
        pdf_reader = PdfReader(uploaded_file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text()
		
        # Split text into chunks
        text_splitter = CharacterTextSplitter(
            separator="\n",
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len
        )
		
        chunks = text_splitter.split_text(text)
        embeddings = OpenAIEmbeddings()
		
        vectorstore = FAISS.from_texts(chunks, embeddings)
		
        if ctx:
            model = ChatOpenAI(model_name="gpt-3.5-turbo-0125", temperature=0.8)
            filled_template = template.replace("{system}", system_prompt)\
                                      .replace("{targetGrade}", targetGrade)\
                                      .replace("{prompt}", user_prompt)
            prompt = ChatPromptTemplate.from_template(template=filled_template)
            retriever = vectorstore.as_retriever(search_type="similarity_score_threshold", search_kwargs={"score_threshold": 0.5, "k": 4})
            
            chain = (
                {"document": retriever, "ctx": RunnablePassthrough()}
                | prompt
                | model
                | StrOutputParser()
            )
            response = chain.invoke(ctx)
            print(response)
            
            apply_bgcolor = WHITE
            if bgcolor == 'black':
                apply_bgcolor = BLACK
            elif bgcolor == 'white':
                apply_bgcolor = WHITE
            elif bgcolor == 'lightblue':
                apply_bgcolor = LIGHTBLUE
            elif bgcolor == 'cream':
                apply_bgcolor = CREAM
            elif bgcolor == 'grey':
                apply_bgcolor = GREY

            apply_fontcolor = BLACK
            if fontcolor == 'black':
                apply_fontcolor = BLACK
            elif fontcolor == 'white':
                apply_fontcolor = WHITE

            file_name = generate_slides_from_XML(response, apply_bgcolor, fonttype, apply_fontcolor)
            return JsonResponse({'filename' : file_name, 'response': response})
        else:
            return JsonResponse({'response': 'No context specified'})
    return HttpResponse("Listening for requests...")

def serve_presentation(request, file_name):
    file_path = os.path.join(GENERATEDCONTENT_DIRECTORY, file_name)
    
    if os.path.exists(file_path) and os.path.isfile(file_path):
        download = request.GET.get('download', 'false').lower() in ['true', '1', 't']
        return FileResponse(open(file_path, 'rb'), as_attachment=download)
    else:
        raise Http404("The requested file does not exist")
    
def regenerate(request):
    return HttpResponse("Listening for requests on regenerate...")